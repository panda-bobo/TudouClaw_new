"""Media tools — create_pptx, create_pptx_advanced, create_video.

All three produce files from structured input, rely on the sandbox for
output-path resolution, and auto-install their heavy third-party deps
(python-pptx, moviepy) on first use.
"""
from __future__ import annotations

import logging
import subprocess
import sys
from typing import Any

from .. import sandbox as _sandbox

logger = logging.getLogger(__name__)


# First-install fallback: when python-pptx / moviepy aren't available,
# we shell out to pip. Cap how long we wait — 60 s is enough for the
# wheel to download on a reasonable connection.
_PIP_INSTALL_TIMEOUT_S = 60


def _ensure_pptx():
    """Import python-pptx, installing it if missing.

    Returns the (Presentation, Inches, Pt, PP_ALIGN) tuple ready to use.
    Raises RuntimeError with the pip stderr if install fails.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        return Presentation, Inches, Pt, PP_ALIGN
    except ImportError:
        result = subprocess.run(
            [sys.executable, "-m", "pip", "install", "python-pptx",
             "--break-system-packages"],
            capture_output=True, text=True, timeout=_PIP_INSTALL_TIMEOUT_S,
        )
        if result.returncode != 0:
            raise RuntimeError(f"Error installing python-pptx: {result.stderr}")
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        return Presentation, Inches, Pt, PP_ALIGN


# ── create_pptx ──────────────────────────────────────────────────────

def _tool_create_pptx(output_path: str, slides: list, title: str = "",
                      **_: Any) -> str:
    """Create a PowerPoint presentation file."""
    try:
        try:
            Presentation, Inches, Pt, PP_ALIGN = _ensure_pptx()
        except RuntimeError as e:
            return str(e)

        pol = _sandbox.get_current_policy()
        output_file = pol.safe_path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()

        # Set title if provided.
        if title:
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title

        # Add slides.
        for slide_data in slides:
            slide_title = slide_data.get("title", "")
            content = slide_data.get("content", "")
            layout_type = slide_data.get("layout", "title_content").lower()

            # Select layout.
            if layout_type == "title":
                layout = prs.slide_layouts[0]
            elif layout_type == "content":
                layout = prs.slide_layouts[5]  # Blank with title
            elif layout_type == "blank":
                layout = prs.slide_layouts[6]  # Blank
            else:  # title_content
                layout = prs.slide_layouts[1]  # Title and content

            slide = prs.slides.add_slide(layout)

            # Add title.
            if slide_title and len(slide.shapes) > 0:
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = slide_title

            # Add content.
            if content and len(slide.shapes) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()

                # Split content by lines and add as bullet points.
                for line in content.split('\n'):
                    if line.strip():
                        p = tf.add_paragraph()
                        p.text = line.strip()
                        p.level = 0

            # Add images onto the slide.
            for img_spec in (slide_data.get("images") or []):
                if not isinstance(img_spec, dict):
                    continue
                img_path_raw = img_spec.get("path", "")
                if not img_path_raw:
                    continue
                img_file = pol.safe_path(img_path_raw)
                left_v = Inches(float(img_spec.get("left", 1)))
                top_v = Inches(float(img_spec.get("top", 2)))
                kw: dict = {}
                if img_spec.get("width"):
                    kw["width"] = Inches(float(img_spec["width"]))
                if img_spec.get("height"):
                    kw["height"] = Inches(float(img_spec["height"]))
                slide.shapes.add_picture(str(img_file), left_v, top_v, **kw)

        prs.save(str(output_file))
        return f"✓ Created presentation: {output_file}"
    except Exception as e:
        return f"Error creating presentation: {e}"


# ── create_pptx_advanced ─────────────────────────────────────────────

def _tool_create_pptx_advanced(
    output_path: str,
    slides: list,
    theme: dict | None = None,
    **_: Any,
) -> str:
    """Create an advanced PowerPoint with shapes, charts, tables, and infographics."""
    try:
        # Auto-install python-pptx (extended imports beyond _ensure_pptx).
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
        except ImportError:
            result = subprocess.run(
                [sys.executable, "-m", "pip", "install", "python-pptx",
                 "--break-system-packages"],
                capture_output=True, text=True,
                timeout=_PIP_INSTALL_TIMEOUT_S,
            )
            if result.returncode != 0:
                return f"Error installing python-pptx: {result.stderr}"
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor

        pol = _sandbox.get_current_policy()
        output_file = pol.safe_path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)

        # Theme defaults.
        th = theme or {}
        T_PRIMARY = th.get("primary", "E8590C")
        T_SECONDARY = th.get("secondary", "2B2B2B")
        T_ACCENT = th.get("accent", "F4A261")
        T_BG = th.get("background", "FFFFFF")
        T_TITLE_FONT = th.get("title_font", "Microsoft YaHei")
        T_BODY_FONT = th.get("body_font", "Microsoft YaHei")

        def _rgb(hex_str: str) -> RGBColor:
            """Convert hex string to RGBColor."""
            h = hex_str.lstrip("#")
            if len(h) == 6:
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            return RGBColor(0, 0, 0)

        def _resolve_color(val: str) -> str:
            """Replace theme placeholders like 'primary' with actual hex."""
            m = {"primary": T_PRIMARY, "secondary": T_SECONDARY,
                 "accent": T_ACCENT, "background": T_BG}
            return m.get(val, val)

        # Shape type mapping.
        SHAPE_MAP = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "arrow_right": MSO_SHAPE.RIGHT_ARROW,
            "arrow_left": MSO_SHAPE.LEFT_ARROW,
            "chevron": MSO_SHAPE.CHEVRON,
            "diamond": MSO_SHAPE.DIAMOND,
            "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "star": MSO_SHAPE.STAR_5_POINT,
        }

        # Chart type mapping.
        CHART_MAP = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "radar": XL_CHART_TYPE.RADAR,
            "area": XL_CHART_TYPE.AREA,
        }

        # Alignment mapping.
        ALIGN_MAP = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }
        VALIGN_MAP = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }

        # Create presentation (16:9 widescreen).
        prs = Presentation()
        SLIDE_W = 10.0   # inches
        SLIDE_H = 5.625  # inches
        MARGIN = 0.15    # minimum margin from edge
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)

        def _clamp_bounds(el: dict) -> dict:
            """Clamp element x/y/w/h so it stays within slide boundaries.

            Fixes the most common LLM layout mistake: placing 3+ items
            in a row where the last one overflows the right/bottom edge.
            """
            x = float(el.get("x", 0))
            y = float(el.get("y", 0))
            w = float(el.get("w", 1))
            h = float(el.get("h", 1))

            # Clamp negative positions.
            if x < 0:
                x = 0
            if y < 0:
                y = 0

            max_w = SLIDE_W - MARGIN
            max_h = SLIDE_H - MARGIN

            # If right edge overflows, shrink width first; if still bad, shift left.
            if x + w > max_w:
                new_w = max_w - x
                if new_w >= w * 0.3 and new_w > 0.3:
                    w = new_w
                else:
                    x = max(0, max_w - w)
                    if x + w > max_w:
                        w = max_w - x

            # Same for bottom edge.
            if y + h > max_h:
                new_h = max_h - y
                if new_h >= h * 0.3 and new_h > 0.3:
                    h = new_h
                else:
                    y = max(0, max_h - h)
                    if y + h > max_h:
                        h = max_h - y

            el["x"] = round(x, 3)
            el["y"] = round(y, 3)
            el["w"] = round(w, 3)
            el["h"] = round(h, 3)
            return el

        def _add_text_element(slide, el):
            """Add a text box to the slide."""
            left = Inches(el.get("x", 0))
            top = Inches(el.get("y", 0))
            width = Inches(el.get("w", 8))
            height = Inches(el.get("h", 1))

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True

            # Background color.
            bg = el.get("bg_color", "")
            if bg:
                fill = txBox.fill
                fill.solid()
                fill.fore_color.rgb = _rgb(_resolve_color(bg))

            # Vertical alignment.
            va = el.get("valign", "")
            if va in VALIGN_MAP:
                tf.paragraphs[0].alignment  # ensure exists
                txBox.text_frame._txBody.attrib  # access
                try:
                    tf._txBody[0].attrib  # bodyPr
                except Exception:
                    pass
                # Set via the text frame directly.
                from pptx.oxml.ns import qn
                bodyPr = tf._txBody.find(qn("a:bodyPr"))
                if bodyPr is not None:
                    anchor_val = {"top": "t", "middle": "ctr",
                                  "bottom": "b"}.get(va, "t")
                    bodyPr.set("anchor", anchor_val)

            content = el.get("content", "")
            lines = content.split("\\n") if "\\n" in content else content.split("\n")

            font_size = el.get("font_size", 14)
            font_name = el.get("font_name", T_BODY_FONT)
            bold = el.get("bold", False)
            italic = el.get("italic", False)
            color = _resolve_color(el.get("color", T_SECONDARY))
            align = el.get("align", "left")
            line_spacing = el.get("line_spacing", 0)

            for i, line_text in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line_text
                p.font.size = Pt(font_size)
                p.font.name = font_name
                p.font.bold = bold
                p.font.italic = italic
                if color:
                    p.font.color.rgb = _rgb(color)
                if align in ALIGN_MAP:
                    p.alignment = ALIGN_MAP[align]
                if line_spacing and line_spacing > 0:
                    p.line_spacing = Pt(font_size * line_spacing)

        def _add_shape_element(slide, el):
            """Add a shape to the slide."""
            shape_type_name = el.get("shape_type", "rectangle")
            mso_shape = SHAPE_MAP.get(shape_type_name, MSO_SHAPE.RECTANGLE)
            left = Inches(el.get("x", 0))
            top = Inches(el.get("y", 0))
            width = Inches(el.get("w", 1))
            height = Inches(el.get("h", 1))

            shape = slide.shapes.add_shape(mso_shape, left, top, width, height)

            fill_color = el.get("fill_color", "")
            if fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = _rgb(_resolve_color(fill_color))
            else:
                shape.fill.background()  # transparent

            line_color = el.get("line_color", "")
            line_width = el.get("line_width", 0)
            if line_color:
                shape.line.color.rgb = _rgb(_resolve_color(line_color))
                shape.line.width = Pt(line_width or 1)
            else:
                shape.line.fill.background()  # no border

            rotation = el.get("rotation", 0)
            if rotation:
                shape.rotation = rotation

        def _add_line_element(slide, el):
            """Add a line connector."""
            x = Inches(el.get("x", 0))
            y = Inches(el.get("y", 0))
            w = Inches(el.get("w", 1))
            h = Inches(el.get("h", 0))

            connector = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR_TYPE.STRAIGHT
                x, y, x + w, y + h,
            )
            lc = el.get("line_color", "CCCCCC")
            if lc:
                connector.line.color.rgb = _rgb(_resolve_color(lc))
            lw = el.get("line_width", 1)
            connector.line.width = Pt(lw)

        def _add_chart_element(slide, el):
            """Add a chart to the slide."""
            from pptx.chart.data import CategoryChartData

            chart_type_name = el.get("chart_type", "column")
            xl_chart = CHART_MAP.get(chart_type_name,
                                     XL_CHART_TYPE.COLUMN_CLUSTERED)

            x = Inches(el.get("x", 0.5))
            y = Inches(el.get("y", 1.5))
            w = Inches(el.get("w", 5))
            h = Inches(el.get("h", 3.5))

            chart_data = CategoryChartData()
            categories = el.get("categories", [])
            chart_data.categories = categories

            for s in (el.get("series") or []):
                chart_data.add_series(
                    s.get("name", "Series"),
                    s.get("values", []),
                )

            chart_frame = slide.shapes.add_chart(xl_chart, x, y, w, h, chart_data)
            chart = chart_frame.chart

            # Apply custom colors.
            colors_list = el.get("colors", [])
            if colors_list:
                try:
                    plot = chart.plots[0]
                    if chart_type_name in ("pie", "doughnut"):
                        # For pie/doughnut, color each point.
                        if plot.series and len(plot.series) > 0:
                            series_obj = plot.series[0]
                            for idx, c in enumerate(colors_list):
                                if idx < len(categories):
                                    point = series_obj.points[idx]
                                    point.format.fill.solid()
                                    point.format.fill.fore_color.rgb = _rgb(_resolve_color(c))
                    else:
                        # For bar/column/line, color each series.
                        for idx, c in enumerate(colors_list):
                            if idx < len(plot.series):
                                s = plot.series[idx]
                                s.format.fill.solid()
                                s.format.fill.fore_color.rgb = _rgb(_resolve_color(c))
                except Exception:
                    pass  # color fail is non-critical

            # Labels.
            if el.get("show_labels") or el.get("show_percent"):
                try:
                    plot = chart.plots[0]
                    plot.has_data_labels = True
                    data_labels = plot.data_labels
                    if el.get("show_percent") and chart_type_name in ("pie", "doughnut"):
                        data_labels.show_percentage = True
                        data_labels.show_value = False
                    else:
                        data_labels.show_value = True
                    data_labels.font.size = Pt(10)
                except Exception:
                    pass

            # Legend.
            if el.get("show_legend") is False:
                chart.has_legend = False
            elif el.get("show_legend"):
                chart.has_legend = True

        def _add_table_element(slide, el):
            """Add a table to the slide."""
            headers = el.get("headers", [])
            rows_data = el.get("rows", [])
            total_rows = len(rows_data) + (1 if headers else 0)
            total_cols = len(headers) if headers else (len(rows_data[0]) if rows_data else 1)

            x = Inches(el.get("x", 0.5))
            y = Inches(el.get("y", 1.5))
            w = Inches(el.get("w", 9))
            h = Inches(el.get("h", 3))

            table_shape = slide.shapes.add_table(total_rows, total_cols, x, y, w, h)
            table = table_shape.table

            # Set column widths evenly.
            col_w = Emu(int(Inches(el.get("w", 9)) / total_cols))
            for ci in range(total_cols):
                table.columns[ci].width = col_w

            header_bg = el.get("header_color", T_PRIMARY)
            header_fc = el.get("header_font_color", "FFFFFF")
            stripe_bg = el.get("stripe_color", "")

            row_offset = 0
            if headers:
                for ci, htext in enumerate(headers):
                    cell = table.cell(0, ci)
                    cell.text = str(htext)
                    # Header style.
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(_resolve_color(header_bg))
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = _rgb(_resolve_color(header_fc))
                        paragraph.font.name = T_BODY_FONT
                        paragraph.alignment = PP_ALIGN.CENTER
                row_offset = 1

            for ri, row in enumerate(rows_data):
                for ci, cval in enumerate(row):
                    if ci >= total_cols:
                        break
                    cell = table.cell(ri + row_offset, ci)
                    cell.text = str(cval)
                    # Stripe.
                    if stripe_bg and ri % 2 == 1:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _rgb(_resolve_color(stripe_bg))
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(11)
                        paragraph.font.name = T_BODY_FONT
                        paragraph.alignment = PP_ALIGN.CENTER

        def _add_icon_circle(slide, el):
            """Add a circle with text inside (for numbering, icons, etc.)."""
            x = Inches(el.get("x", 0))
            y = Inches(el.get("y", 0))
            w = Inches(el.get("w", 0.8))
            h = Inches(el.get("h", 0.8))

            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
            fill_color = el.get("fill_color", T_PRIMARY)
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(_resolve_color(fill_color))
            shape.line.fill.background()  # no border

            # Add text inside.
            tf = shape.text_frame
            tf.word_wrap = False
            from pptx.oxml.ns import qn
            bodyPr = tf._txBody.find(qn("a:bodyPr"))
            if bodyPr is not None:
                bodyPr.set("anchor", "ctr")

            p = tf.paragraphs[0]
            p.text = el.get("text", "")
            p.font.size = Pt(el.get("font_size", 16))
            p.font.bold = True
            p.font.color.rgb = _rgb(_resolve_color(el.get("font_color", "FFFFFF")))
            p.font.name = T_TITLE_FONT
            p.alignment = PP_ALIGN.CENTER

        def _add_image_element(slide, el):
            """Add an image to the slide."""
            img_path = el.get("path", "")
            if not img_path:
                return
            img_file = pol.safe_path(img_path)
            if not img_file.exists():
                return
            x = Inches(el.get("x", 0))
            y = Inches(el.get("y", 0))
            kw = {}
            if el.get("w"):
                kw["width"] = Inches(el["w"])
            if el.get("h"):
                kw["height"] = Inches(el["h"])
            slide.shapes.add_picture(str(img_file), x, y, **kw)

        # Element dispatcher.
        ELEMENT_HANDLERS = {
            "text": _add_text_element,
            "shape": _add_shape_element,
            "line": _add_line_element,
            "chart": _add_chart_element,
            "table": _add_table_element,
            "icon_circle": _add_icon_circle,
            "image": _add_image_element,
        }

        # Layout engine integration.
        try:
            from ..utils.pptx_layouts import generate_layout
        except ImportError:
            generate_layout = None

        # Build slides.
        blank_layout = prs.slide_layouts[6]  # Blank
        slide_count = 0
        theme_dict = {
            "primary": T_PRIMARY, "secondary": T_SECONDARY,
            "accent": T_ACCENT, "background": T_BG,
            "title_font": T_TITLE_FONT, "body_font": T_BODY_FONT,
        }

        for slide_data in slides:
            slide = prs.slides.add_slide(blank_layout)
            slide_count += 1

            # Slide background.
            bg_color = slide_data.get("background", T_BG)
            if bg_color:
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = _rgb(_resolve_color(bg_color))

            # If layout spec is present, auto-generate elements from it.
            all_elements = []
            layout_spec = slide_data.get("layout")
            if layout_spec and generate_layout:
                try:
                    auto_els = generate_layout(layout_spec, theme_dict)
                    all_elements.extend(auto_els)
                except Exception as _le:
                    logger.warning("pptx_advanced layout error: %s", _le)

            # Append any manually-specified elements (can supplement layout).
            all_elements.extend(slide_data.get("elements") or [])

            # Add elements in order (z-order: first = bottom).
            for el in all_elements:
                el_type = el.get("type", "")
                handler = ELEMENT_HANDLERS.get(el_type)
                if handler:
                    try:
                        # Clamp bounds to prevent overflow.
                        if any(k in el for k in ("x", "y", "w", "h")):
                            _clamp_bounds(el)
                        handler(slide, el)
                    except Exception as e:
                        # Non-critical: log but continue.
                        logger.warning("pptx_advanced element error (%s): %s",
                                       el_type, e)

        prs.save(str(output_file))
        return f"✓ Created advanced presentation ({slide_count} slides): {output_file}"

    except Exception as e:
        return f"Error creating presentation: {e}"


# ── create_video ─────────────────────────────────────────────────────

def _tool_create_video(output_path: str, frames: list, fps: int = 24,
                      audio_path: str = "", **_: Any) -> str:
    """Create a video from image frames."""
    try:
        pol = _sandbox.get_current_policy()
        output_file = pol.safe_path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)

        # Validate frame paths.
        frame_list = []
        for frame in frames:
            img_path = frame.get("image_path", "")
            if not img_path:
                return "Error: Each frame must have image_path"
            img_file = pol.safe_path(img_path)
            if not img_file.exists():
                return f"Error: Image file not found: {img_path}"
            duration = frame.get("duration", 3)
            frame_list.append((str(img_file), duration))

        # Try moviepy first.
        try:
            import moviepy.editor as mpy
        except ImportError:
            result = subprocess.run(
                [sys.executable, "-m", "pip", "install", "moviepy",
                 "--break-system-packages"],
                capture_output=True, text=True,
                timeout=_PIP_INSTALL_TIMEOUT_S,
            )
            if result.returncode != 0:
                return f"Error installing moviepy: {result.stderr}"
            import moviepy.editor as mpy

        # Create video from frames.
        clips = []
        for img_path, duration in frame_list:
            clip = mpy.ImageClip(img_path).set_duration(duration)
            clips.append(clip)

        video = mpy.concatenate_videoclips(clips)

        # Add audio if provided.
        if audio_path:
            audio_file = pol.safe_path(audio_path)
            if audio_file.exists():
                audio = mpy.AudioFileClip(str(audio_file))
                video = video.set_audio(audio)

        # Write video file.
        video.write_videofile(str(output_file), fps=fps, verbose=False,
                              logger=None)
        video.close()

        return f"✓ Video created: {output_path}"
    except Exception as e:
        return f"Error creating video: {e}"
